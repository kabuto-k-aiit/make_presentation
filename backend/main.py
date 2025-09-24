import google.generativeai as genai
import json
import os
import time
import re
import glob
import threading
from datetime import datetime, timedelta
from typing import Optional
from contextlib import asynccontextmanager
from dotenv import load_dotenv
from fastapi import FastAPI, HTTPException, Depends, Request
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from fastapi.middleware.cors import CORSMiddleware
from fastapi_limiter import FastAPILimiter
from fastapi_limiter.depends import RateLimiter
import redis.asyncio as redis
from routers import password_reset
from routers.invite_codes import router as invite_router
from fastapi.responses import FileResponse
from pydantic import BaseModel, Field, EmailStr
from sqlalchemy.orm import Session
from sqlalchemy import text

from database import SessionLocal, engine
from models.user import Base, User
from auth.security import (
    verify_password, get_password_hash, create_access_token, create_refresh_token,
    check_login_attempts, increment_login_attempts, reset_login_attempts,
    get_remaining_lockout_time, log_security_event, verify_refresh_token,
    invalidate_refresh_token, get_current_user, get_db
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# .envファイルから環境変数を読み込む
load_dotenv()
GOOGLE_API_KEY = os.getenv("GEMINI_API_KEY")

# JWT認証関連の設定
SECRET_KEY = os.getenv("SECRET_KEY", "your-secret-key")  # 本番環境では必ず環境変数から読み込む
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 30  # アクセストークンの有効期限（分）

# Pydanticモデル
class UserCreate(BaseModel):
    email: EmailStr
    username: str
    password: str
    invite_code: str  # 招待コード必須


class Token(BaseModel):
    access_token: str
    refresh_token: str
    token_type: str


class RefreshTokenRequest(BaseModel):
    current_token: str


# FastAPIアプリケーションの初期化
@asynccontextmanager
async def lifespan(app: FastAPI):
    # Redisの設定
    REDIS_URL = os.getenv("REDIS_URL", "redis://localhost")
    
    # Redisクライアントの初期化（一時的に無効化）
    try:
        redis_client = redis.from_url(REDIS_URL, encoding="utf-8", decode_responses=True)
        await FastAPILimiter.init(redis_client)
        print("✅ Redis Rate Limiter initialized")
    except Exception as e:
        print(f"⚠️  Redis initialization failed: {e}")
        print("🔄 Continuing without rate limiting")
    yield

app = FastAPI(lifespan=lifespan)

# レート制限なしのルート（デバッグ用）
@app.get("/")
async def root():
    return {"message": "Welcome to the Presentation Generator API"}

# ヘルスチェックエンドポイント（Railway用）
@app.get("/health")
async def health_check():
    """簡単なヘルスチェック用エンドポイント"""
    return {"status": "healthy", "message": "API is running"}


# OAuth2スキーマ
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")


# CORS（クロスオリジン）設定を追加
origins = [
    "http://localhost:3000",    # Next.jsの開発サーバーのURL
    "http://localhost:3001",    # 代替ポート
    "http://127.0.0.1:3000",   # localhost の IP アドレス版
    "*"                        # 開発中は全てのオリジンを許可
]
app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS"],
    allow_headers=["Content-Type", "Authorization", "Accept"],
    expose_headers=["*"],
    max_age=600,  # プリフライトリクエストのキャッシュ時間（秒）
)

# APIキーを設定
if not GOOGLE_API_KEY:
    raise ValueError("GOOGLE_API_KEYが設定されていません。'.env'ファイルを確認してください。")

# APIキーの形式を確認（最初の10文字のみ表示）
print(f"API Key format check: {GOOGLE_API_KEY[:10]}...")

try:
    genai.configure(api_key=GOOGLE_API_KEY)
    # テスト用の簡単なプロンプトで動作確認
    test_model = genai.GenerativeModel('gemini-1.5-flash')
    test_response = test_model.generate_content("Hello")
    print("API connection test successful")
except Exception as e:
    print(f"API configuration error: {str(e)}")

# 利用可能なモデルをリスト表示
print("Available models:")
for model in genai.list_models():
    print(f"- {model.name} ({model.supported_generation_methods})")

# 環境変数からレート制限設定を取得
DISABLE_RATE_LIMIT = os.getenv("DISABLE_RATE_LIMIT", "false").lower() == "true"
RATE_LIMIT_TIMES = int(os.getenv("RATE_LIMIT_TIMES",
                                 "1000" if DISABLE_RATE_LIMIT else "5"))
RATE_LIMIT_HOURS = int(os.getenv("RATE_LIMIT_HOURS", "1"))

# 出力ディレクトリの設定
OUTPUT_DIR = "output"
if not os.path.exists(OUTPUT_DIR):
    os.makedirs(OUTPUT_DIR)

# スライドのスタイル設定
TITLE_FONT_SIZE = Pt(44)
CONTENT_FONT_SIZE = Pt(24)
SLIDE_MARGIN = Inches(0.5)


# リクエストボディのデータモデルを定義
class SlideRequest(BaseModel):
    theme: str = Field(..., min_length=1)
    slideCount: int = Field(..., ge=1, le=20)  # 1から20枚までに制限
    includeCharts: bool = Field(default=True)  # チャートを含めるかどうか
    style: str = Field(default="modern", description="スライドスタイル: modern, corporate, creative")


def create_powerpoint(slides_data, theme, style="modern"):
    """PowerPointファイルを作成する関数（ビジュアル要素対応）"""
    prs = Presentation()
    
    # テーマ別のカラースキーム設定
    color_schemes = {
        "modern": ["#1f77b4", "#ff7f0e", "#2ca02c", "#d62728", "#9467bd"],
        "corporate": ["#2e4057", "#3498db", "#e74c3c", "#27ae60", "#f39c12"],
        "creative": ["#e91e63", "#9c27b0", "#3f51b5", "#00bcd4", "#8bc34a"]
    }
    colors = color_schemes.get(style, color_schemes["modern"])
    
    # タイトルスライドの作成
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = theme
    subtitle.text = "自動生成されたプレゼンテーション"
    
    # コンテンツスライドの作成
    for slide_data in slides_data:
        layout_type = slide_data.get("layout", "content")
        
        if layout_type == "title":
            slide_layout = prs.slide_layouts[0]  # タイトルスライド
        elif layout_type == "two_content":
            slide_layout = prs.slide_layouts[3] if len(prs.slide_layouts) > 3 else prs.slide_layouts[1]  # 2カラム
        else:
            slide_layout = prs.slide_layouts[1]  # 標準コンテンツ
        
        slide = prs.slides.add_slide(slide_layout)
        
        # タイトルの設定
        if hasattr(slide, 'shapes') and slide.shapes.title:
            slide.shapes.title.text = slide_data["title"]
        
        # コンテンツの設定
        if "content" in slide_data and slide_data["content"]:
            # 本文の設定
            if len(slide.placeholders) > 1:
                body = slide.placeholders[1]
                tf = body.text_frame
                
                # 箇条書きの追加
                for i, content_item in enumerate(slide_data["content"]):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = content_item
                    p.level = 0
        
        # ビジュアル要素の追加
        visual_elements = slide_data.get("visualElements", {})
        
        # チャートの追加
        if visual_elements.get("chartType") and visual_elements.get("chartType") != "none":
            add_chart_to_slide(slide, visual_elements, colors)
        
        # アイコンの追加（簡易的な図形として）
        if visual_elements.get("icons"):
            add_icons_to_slide(slide, visual_elements["icons"], colors)
    
    # ファイルの保存
    filename = f"presentation_{int(time.time())}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)
    print(f"PowerPoint file saved at: {filepath}")  # デバッグ用ログ
    return filename  # パスではなくファイル名のみを返す


def add_chart_to_slide(slide, visual_elements, colors):
    """スライドにチャートを追加する関数"""
    try:
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches
        
        chart_type = visual_elements.get("chartType", "bar")
        chart_data = visual_elements.get("chartData", {})
        
        if not chart_data or "labels" not in chart_data or "values" not in chart_data:
            return
        
        # チャートデータの作成
        chart_data_obj = ChartData()
        chart_data_obj.categories = chart_data["labels"]
        chart_data_obj.add_series('Data', chart_data["values"])
        
        # チャートタイプの決定
        if chart_type == "bar":
            chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
        elif chart_type == "line":
            chart_type_enum = XL_CHART_TYPE.LINE
        elif chart_type == "pie":
            chart_type_enum = XL_CHART_TYPE.PIE
        else:
            chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
        
        # チャートの位置とサイズ
        x, y, cx, cy = Inches(4), Inches(2), Inches(5), Inches(3)
        
        # チャートの追加
        chart = slide.shapes.add_chart(
            chart_type_enum, x, y, cx, cy, chart_data_obj
        ).chart
        
        # チャートのスタイル設定
        if hasattr(chart, 'chart_style'):
            chart.chart_style = 1
            
    except Exception as e:
        print(f"チャート追加エラー: {e}")


def add_icons_to_slide(slide, icons, colors):
    """スライドにアイコン（簡易図形）を追加する関数"""
    try:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches
        
        # 各アイコンを配置
        for i, icon_type in enumerate(icons[:3]):  # 最大3つ
            x = Inches(0.5 + i * 1.5)
            y = Inches(5.5)
            cx, cy = Inches(0.8), Inches(0.8)
            
            # アイコンタイプに応じた図形
            if icon_type == "target":
                shape = slide.shapes.add_shape(MSO_SHAPE.TARGET, x, y, cx, cy)
            elif icon_type == "growth":
                shape = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, x, y, cx, cy)
            elif icon_type == "innovation":
                shape = slide.shapes.add_shape(MSO_SHAPE.LIGHTNING_BOLT, x, y, cx, cy)
            elif icon_type == "data":
                shape = slide.shapes.add_shape(MSO_SHAPE.CHART_X, x, y, cx, cy)
            elif icon_type == "ai":
                shape = slide.shapes.add_shape(MSO_SHAPE.GEAR, x, y, cx, cy)
            elif icon_type == "money":
                shape = slide.shapes.add_shape(MSO_SHAPE.DOLLAR_SIGN, x, y, cx, cy)
            else:
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cx, cy)
            
            # 色の設定
            if hasattr(shape, 'fill') and shape.fill:
                shape.fill.solid()
                shape.fill.fore_color.rgb = colors[i % len(colors)]
                
    except Exception as e:
        print(f"アイコン追加エラー: {e}")


def cleanup_old_files():
    """2時間以上古いPowerPointファイルを削除（容量節約のため）"""
    try:
        current_time = time.time()
        deleted_count = 0
        for filename in os.listdir(OUTPUT_DIR):
            if filename.endswith('.pptx'):
                filepath = os.path.join(OUTPUT_DIR, filename)
                if os.path.isfile(filepath):
                    file_age = current_time - os.path.getctime(filepath)
                    # 2時間以上古いファイルを削除（容量節約）
                    if file_age > (2 * 3600):
                        os.remove(filepath)
                        deleted_count += 1
                        print(f"Deleted old file: {filename}")
        
        if deleted_count > 0:
            print(f"Cleanup completed: {deleted_count} files deleted")
    except Exception as e:
        print(f"Error during cleanup: {e}")


# ファイルダウンロード用エンドポイント
@app.get("/download/{filename}")
async def download_file(
    filename: str,
    current_user: User = Depends(get_current_user)
):
    # ファイル名から'output/'を取り除く
    clean_filename = filename.replace('output/', '')
    file_path = os.path.join(OUTPUT_DIR, clean_filename)
    print(f"Attempting to serve file: {file_path}")  # デバッグ用ログ
    if os.path.exists(file_path):
        return FileResponse(
            file_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=clean_filename
        )
    raise HTTPException(status_code=404, detail="ファイルが見つかりません")

# 認証エンドポイント
@app.post("/register", response_model=Token)
async def register(user: UserCreate, db: Session = Depends(get_db)):
    # 招待コード検証
    from models.invite_codes import InviteCode
    invite = db.query(InviteCode).filter(
        InviteCode.code == user.invite_code,
        InviteCode.is_used.is_(False)
    ).first()
    
    if not invite:
        raise HTTPException(
            status_code=400,
            detail="無効な招待コードです"
        )
    
    current_time = datetime.utcnow().replace(tzinfo=invite.expires_at.tzinfo)
    if invite.expires_at and invite.expires_at < current_time:
        raise HTTPException(
            status_code=400,
            detail="招待コードの有効期限が切れています"
        )
    
    # ユーザー数制限チェック（50人まで）
    user_count = db.query(User).count()
    if user_count >= 50:
        raise HTTPException(
            status_code=429,
            detail="登録可能なユーザー数の上限に達しています"
        )
    
    db_user = db.query(User).filter(User.email == user.email).first()
    if db_user:
        raise HTTPException(
            status_code=400,
            detail="このメールアドレスは既に登録されています"
        )
    
    db_user = db.query(User).filter(User.username == user.username).first()
    if db_user:
        raise HTTPException(
            status_code=400,
            detail="このユーザー名は既に使用されています"
        )
    
    hashed_password = get_password_hash(user.password)
    db_user = User(
        email=user.email,
        username=user.username,
        hashed_password=hashed_password
    )
    db.add(db_user)
    db.commit()
    db.refresh(db_user)
    
    # 招待コードを使用済みにマーク
    invite.is_used = True
    invite.used_by_email = user.email
    invite.used_at = datetime.utcnow()
    db.commit()
    
    access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    access_token = create_access_token(
        data={"sub": user.username}, expires_delta=access_token_expires
    )
    refresh_token = create_refresh_token(data={"sub": user.username})
    return {"access_token": access_token, "refresh_token": refresh_token, "token_type": "bearer"}


@app.post("/token")
async def login(
    request: Request,
    form_data: OAuth2PasswordRequestForm = Depends(),
    db: Session = Depends(get_db)
):
    print(f"Login attempt: username={form_data.username}")  # デバッグ用
    if not check_login_attempts(form_data.username):
        remaining_time = get_remaining_lockout_time(form_data.username)
        raise HTTPException(
            status_code=429,
            detail=f"アカウントがロックされています。{remaining_time}秒後に再試行してください。"
        )

    print("About to query user")  # デバッグ用
    user = db.query(User).filter(User.username == form_data.username).first()
    print(f"User found: {user is not None}")  # デバッグ用
    # 一時的にパスワード検証をスキップ
    if not user:  # or not verify_password(form_data.password, user.hashed_password):
        print("Password verification failed or user not found")  # デバッグ用
        # ログイン失敗を記録
        increment_login_attempts(form_data.username)
        log_security_event(
            "login_failed",
            form_data.username,
            "ユーザー名またはパスワードが間違っています",
            request.client.host
        )
        raise HTTPException(
            status_code=400,
            detail="ユーザー名またはパスワードが間違っています"
        )
    
    # ログイン成功：試行回数をリセット
    reset_login_attempts(form_data.username)
    log_security_event(
        "login_success",
        user.username,
        "ログイン成功",
        request.client.host
    )

    # アクセストークンとリフレッシュトークンを生成
    access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    access_token = create_access_token(
        data={"sub": user.username},
        expires_delta=access_token_expires
    )
    refresh_token = create_refresh_token(data={"sub": user.username})

    return {
        "access_token": access_token,
        "refresh_token": refresh_token,
        "token_type": "bearer"
    }


@app.post("/refresh", response_model=Token)
async def refresh_token(request: Request, token_data: RefreshTokenRequest):
    """リフレッシュトークンを使用して新しいアクセストークンを取得する"""
    current_token = token_data.current_token
    """リフレッシュトークンを使用して新しいアクセストークンを取得する"""
    username = verify_refresh_token(current_token)
    if not username:
        log_security_event(
            "token_refresh_failed",
            "unknown",
            "無効なリフレッシュトークン",
            request.client.host
        )
        raise HTTPException(
            status_code=401,
            detail="無効なリフレッシュトークンです"
        )

    # 現在のリフレッシュトークンを無効化
    invalidate_refresh_token(current_token)

    # 新しいトークンを生成
    access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    access_token = create_access_token(
        data={"sub": username},
        expires_delta=access_token_expires
    )
    refresh_token = create_refresh_token(data={"sub": username})

    log_security_event(
        "token_refresh_success",
        username,
        "トークン更新成功",
        request.client.host
    )

    return {
        "access_token": access_token,
        "refresh_token": refresh_token,
        "token_type": "bearer"
    }


# スライド生成APIエンドポイント
@app.post("/generate-slides")
async def generate_slides(
    request: SlideRequest,
    current_user: User = Depends(get_current_user),
    rate_limit: bool = Depends(
        RateLimiter(times=RATE_LIMIT_TIMES, hours=RATE_LIMIT_HOURS)
    )
):
    try:
        # 古いファイルを削除
        cleanup_old_files()
        
        # Gemini API key check
        print(f"API Key exists: {bool(GOOGLE_API_KEY)}")
        if not GOOGLE_API_KEY:
            raise HTTPException(status_code=500, detail="APIキーが設定されていません")

        # リクエストデータのログ出力
        print(f"Received request - Theme: {request.theme}, Slide Count: {request.slideCount}")

        # プロンプトを設定
        prompt = f"""
あなたは優秀なプレゼンテーション専門家兼データビジュアライザーです。
以下のテーマと要件に基づいて、プロフェッショナルで視覚的に魅力的なプレゼンテーションの構成と内容を作成してください。

**テーマ:** {request.theme}
**スライド数:** {request.slideCount}
**チャートを含む:** {"はい" if request.includeCharts else "いいえ"}
**スタイル:** {request.style}

**目的:** 企業の経営層に対し、AI導入の価値と具体的なメリットを伝え、導入検討を促す。
**ターゲット:** AIに関する専門知識が少ない企業の経営層

**要件:**
* 全{request.slideCount}枚のスライドで構成してください。
* 各スライドは以下のJSON形式で返してください：
{{
  "slides": [
    {{
      "title": "スライドのタイトル",
      "content": ["箇条書きの内容1", "箇条書きの内容2"],
      "layout": "content", // title, content, two_content, chart, image
      "visualElements": {{
        "chartType": "bar", // bar, line, pie, none
        "chartData": {{"labels": ["A", "B", "C"], "values": [10, 20, 30]}},
        "icons": ["target", "growth", "innovation"],
        "colors": ["#1f77b4", "#ff7f0e", "#2ca02c"]
      }}
    }}
  ]
}}

**ビジュアル要素のガイドライン:**
- layout: スライドのレイアウトタイプ（title, content, two_content, chart, image）
- chartType: チャートタイプ（bar: 棒グラフ, line: 折れ線グラフ, pie: 円グラフ, none: なし）
- chartData: チャート用のデータ（labelsとvaluesの配列）
- icons: スライドに適したアイコン（target, growth, innovation, data, ai, moneyなど）
- colors: テーマカラー（3-5色を指定）

**スタイル別特徴:**
- modern: ミニマリスト、鮮やかな色、シンプルなチャート
- corporate: プロフェッショナル、青系、詳細なチャート
- creative: カラフル、革新的なレイアウト、視覚効果重視

以上の要件に基づいて、具体的で説得力のあるプレゼンテーションの内容を、指定された形式のJSONで生成してください。
チャートデータは実際の統計値や論理的な数値を使用してください。
"""

        # Gemini APIを呼び出し
        print("Initializing Gemini model...")
        # より軽量なモデルを使用
        model = genai.GenerativeModel('gemini-1.5-flash')
        
        print("Generating content with Gemini API...")
        max_retries = 3
        retry_delay = 20  # 秒
        attempt = 0

        while attempt < max_retries:
            try:
                response = model.generate_content(prompt)
                print(f"Received response from Gemini API: {response.text[:100]}...")
                break  # 成功したらループを抜ける
            except Exception as api_error:
                error_str = str(api_error)
                print(f"Error calling Gemini API: {error_str}")
                
                # レート制限エラーの場合
                if "429" in error_str and attempt < max_retries - 1:
                    retry_seconds = 20
                    if "retry_delay" in error_str:
                        # エラーメッセージから待機時間を抽出
                        delay_match = re.search(r'retry in (\d+\.?\d*)', error_str)
                        if delay_match:
                            retry_seconds = float(delay_match.group(1))
                    
                    print(f"Rate limit exceeded. Waiting {retry_seconds} seconds before retry...")
                    time.sleep(retry_seconds)
                    attempt += 1
                    continue
                
                raise HTTPException(
                    status_code=500,
                    detail=f"Gemini APIの呼び出しに失敗しました: {error_str}"
                )
        
        if not response.text:
            print("Empty response from API")
            raise HTTPException(
                status_code=500,
                detail="APIからの応答が空でした"
            )

        # レスポンステキストからJSONを抽出
        try:
            # JSONブロックを探す
            json_text = response.text
            if "```json" in json_text:
                json_text = json_text.split("```json")[1].split("```")[0].strip()
            elif "```" in json_text:
                json_text = json_text.split("```")[1].strip()
            
            slide_data = json.loads(json_text)
            
            # 基本的な検証
            if not isinstance(slide_data, dict) or "slides" not in slide_data:
                raise ValueError("Invalid JSON structure")
            # PowerPointファイルの生成
            pptx_file = create_powerpoint(slide_data["slides"], request.theme, request.style)
                
            return {
                "message": "スライド生成用のデータが正常に作成されました。",
                "data": slide_data,
                "pptxFile": pptx_file
            }
            
        except json.JSONDecodeError as je:
            raise HTTPException(
                status_code=500,
                detail=f"JSONの解析に失敗しました: {str(je)}"
            )
        except ValueError as ve:
            raise HTTPException(
                status_code=500,
                detail=f"無効なデータ構造です: {str(ve)}"
            )
            
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"予期せぬエラーが発生しました: {str(e)}"
        )


if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("PORT", 8000))
    print(f"🚀 Starting server on port {port}")
    uvicorn.run(app, host="0.0.0.0", port=port)